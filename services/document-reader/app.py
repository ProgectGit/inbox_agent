import base64
import io
import json
import re
import zipfile
from datetime import date, datetime, time
from http.server import BaseHTTPRequestHandler, ThreadingHTTPServer
from pathlib import Path
from xml.etree import ElementTree

from docx import Document
from openpyxl import load_workbook
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


MAX_BODY_BYTES = 28 * 1024 * 1024
MAX_FILE_BYTES = 20 * 1024 * 1024
MAX_UNCOMPRESSED_BYTES = 150 * 1024 * 1024
MAX_ZIP_ENTRIES = 10_000
MAX_OUTPUT_CHARS = 120_000
MAX_SHEETS = 100
MAX_ROWS_PER_SHEET = 10_000
MAX_CELLS = 100_000
MAX_SLIDES = 500

SUPPORTED_EXTENSIONS = {".docx", ".xlsx", ".pptx"}
WORD_NS = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"


def clean_text(value) -> str:
    text = str(value or "").replace("\x00", "")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r" *\n *", "\n", text)
    return re.sub(r"\n{3,}", "\n\n", text).strip()


def json_value(value) -> str:
    if value is None:
        return ""
    if isinstance(value, (datetime, date, time)):
        return value.isoformat()
    if isinstance(value, bool):
        return "TRUE" if value else "FALSE"
    return clean_text(value)


def validate_office_zip(data: bytes):
    try:
        archive = zipfile.ZipFile(io.BytesIO(data))
    except zipfile.BadZipFile as exc:
        raise ValueError("invalid_office_archive") from exc
    infos = archive.infolist()
    if len(infos) > MAX_ZIP_ENTRIES:
        archive.close()
        raise ValueError("too_many_archive_entries")
    if sum(item.file_size for item in infos) > MAX_UNCOMPRESSED_BYTES:
        archive.close()
        raise ValueError("office_archive_too_large")
    return archive


def docx_xml_text(archive: zipfile.ZipFile, member: str) -> str:
    if member not in archive.namelist():
        return ""
    try:
        root = ElementTree.fromstring(archive.read(member))
    except ElementTree.ParseError:
        return ""
    paragraphs = []
    for paragraph in root.iter(f"{WORD_NS}p"):
        fragments = [node.text or "" for node in paragraph.iter(f"{WORD_NS}t")]
        value = clean_text("".join(fragments))
        if value:
            paragraphs.append(value)
    return "\n".join(paragraphs)


def docx_textbox_text(archive: zipfile.ZipFile) -> str:
    member = "word/document.xml"
    if member not in archive.namelist():
        return ""
    try:
        root = ElementTree.fromstring(archive.read(member))
    except ElementTree.ParseError:
        return ""
    values = []
    for container in root.iter(f"{WORD_NS}txbxContent"):
        for paragraph in container.iter(f"{WORD_NS}p"):
            fragments = [node.text or "" for node in paragraph.iter(f"{WORD_NS}t")]
            value = clean_text("".join(fragments))
            if value:
                values.append(value)
    return "\n".join(values)


def extract_docx(data: bytes, archive: zipfile.ZipFile) -> dict:
    document = Document(io.BytesIO(data))
    sections = []

    core = document.core_properties
    properties = {
        "title": clean_text(core.title),
        "subject": clean_text(core.subject),
        "author": clean_text(core.author),
        "keywords": clean_text(core.keywords),
        "comments": clean_text(core.comments),
        "created": core.created.isoformat() if core.created else None,
        "modified": core.modified.isoformat() if core.modified else None,
    }
    header = [
        f"Назва: {properties['title']}" if properties["title"] else "",
        f"Автор: {properties['author']}" if properties["author"] else "",
        f"Тема: {properties['subject']}" if properties["subject"] else "",
        f"Ключові слова: {properties['keywords']}" if properties["keywords"] else "",
    ]
    if any(header):
        sections.append("Властивості документа:\n" + "\n".join(item for item in header if item))

    paragraphs = []
    for paragraph in document.paragraphs:
        text = clean_text(paragraph.text)
        if not text:
            continue
        style = clean_text(getattr(paragraph.style, "name", ""))
        if style.lower().startswith("heading"):
            text = f"## {text}"
        paragraphs.append(text)
    if paragraphs:
        sections.append("Основний текст:\n" + "\n".join(paragraphs))

    textbox_text = clean_text(docx_textbox_text(archive))
    if textbox_text:
        sections.append("Текстові блоки:\n" + textbox_text)

    table_sections = []
    for table_index, table in enumerate(document.tables, start=1):
        rows = []
        for row in table.rows:
            values = [clean_text(cell.text).replace("\n", " / ") for cell in row.cells]
            if any(values):
                rows.append(" | ".join(values))
        if rows:
            table_sections.append(f"Таблиця {table_index}:\n" + "\n".join(rows))
    if table_sections:
        sections.append("\n\n".join(table_sections))

    headers = []
    footers = []
    for section in document.sections:
        header_text = clean_text("\n".join(p.text for p in section.header.paragraphs))
        footer_text = clean_text("\n".join(p.text for p in section.footer.paragraphs))
        if header_text and header_text not in headers:
            headers.append(header_text)
        if footer_text and footer_text not in footers:
            footers.append(footer_text)
    if headers:
        sections.append("Колонтитули зверху:\n" + "\n".join(headers))
    if footers:
        sections.append("Колонтитули знизу:\n" + "\n".join(footers))

    extras = []
    for label, member in (
        ("Виноски", "word/footnotes.xml"),
        ("Кінцеві примітки", "word/endnotes.xml"),
        ("Коментарі", "word/comments.xml"),
    ):
        value = clean_text(docx_xml_text(archive, member))
        if value:
            extras.append(f"{label}:\n{value}")
    if extras:
        sections.append("\n\n".join(extras))

    image_count = sum(1 for name in archive.namelist() if name.startswith("word/media/"))
    content = clean_text("\n\n".join(sections))
    return {
        "content": content,
        "metadata": {
            "format": "docx",
            "paragraphs": len(paragraphs),
            "tables": len(document.tables),
            "sections": len(document.sections),
            "images": image_count,
            "properties": properties,
        },
        "needs_ai_fallback": len(content) < 120 and image_count > 0,
    }


def extract_xlsx(data: bytes, _archive: zipfile.ZipFile) -> dict:
    workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=False, keep_links=False)
    sheets_total = len(workbook.sheetnames)
    output = []
    cells_read = 0
    truncated = False
    sheet_metadata = []

    for sheet_index, sheet in enumerate(workbook.worksheets[:MAX_SHEETS], start=1):
        sheet_lines = []
        rows_read = 0
        nonempty_cells = 0
        for row_index, row in enumerate(sheet.iter_rows(), start=1):
            if row_index > MAX_ROWS_PER_SHEET or cells_read >= MAX_CELLS:
                truncated = True
                break
            values = []
            for cell in row:
                value = json_value(cell.value)
                if not value:
                    continue
                cells_read += 1
                nonempty_cells += 1
                values.append(f"{get_column_letter(cell.column)}{cell.row}={value}")
                if cells_read >= MAX_CELLS:
                    truncated = True
                    break
            if values:
                rows_read += 1
                sheet_lines.append(" | ".join(values))
            if sum(len(part) for part in output) + sum(len(part) for part in sheet_lines) >= MAX_OUTPUT_CHARS:
                truncated = True
                break
        sheet_metadata.append({
            "name": sheet.title,
            "rows_read": rows_read,
            "nonempty_cells": nonempty_cells,
            "max_row": sheet.max_row,
            "max_column": sheet.max_column,
        })
        if sheet_lines:
            output.append(f"Аркуш {sheet_index}: {sheet.title}\n" + "\n".join(sheet_lines))
        if truncated:
            break
    if sheets_total > MAX_SHEETS:
        truncated = True
    workbook.close()
    content = clean_text("\n\n".join(output))
    return {
        "content": content,
        "metadata": {
            "format": "xlsx",
            "sheets_total": sheets_total,
            "sheets_read": len(sheet_metadata),
            "cells_read": cells_read,
            "sheets": sheet_metadata,
            "truncated": truncated,
        },
        "needs_ai_fallback": len(content) < 40,
    }


def extract_pptx(data: bytes, archive: zipfile.ZipFile) -> dict:
    presentation = Presentation(io.BytesIO(data))
    slides = []
    image_count = 0
    table_count = 0
    chart_count = 0
    truncated = False

    def read_shapes(shapes, parts, slide_title):
        nonlocal image_count, table_count, chart_count
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                read_shapes(shape.shapes, parts, slide_title)
                continue
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_count += 1
                try:
                    props = shape._element.xpath(".//*[local-name()='cNvPr']")
                    alt = clean_text((props[0].get("descr") or props[0].get("title")) if props else "")
                except (AttributeError, IndexError):
                    alt = ""
                if alt:
                    parts.append("Опис зображення: " + alt)
            if getattr(shape, "has_text_frame", False):
                text = clean_text(shape.text)
                if text and text != slide_title:
                    parts.append(text)
            if getattr(shape, "has_table", False):
                table_count += 1
                rows = []
                for row in shape.table.rows:
                    values = [clean_text(cell.text).replace("\n", " / ") for cell in row.cells]
                    if any(values):
                        rows.append(" | ".join(values))
                if rows:
                    parts.append("Таблиця:\n" + "\n".join(rows))
            if getattr(shape, "has_chart", False):
                chart_count += 1
                chart = shape.chart
                chart_parts = []
                try:
                    if chart.has_title:
                        chart_parts.append("Назва: " + clean_text(chart.chart_title.text_frame.text))
                except (AttributeError, ValueError):
                    pass
                for series in chart.series:
                    name = clean_text(series.name)
                    values = []
                    try:
                        values = [json_value(value) for value in series.values]
                    except (AttributeError, TypeError, ValueError):
                        pass
                    line = f"{name}: {', '.join(value for value in values if value)}".strip(": ")
                    if line:
                        chart_parts.append(line)
                if chart_parts:
                    parts.append("Діаграма:\n" + "\n".join(chart_parts))

    for slide_index, slide in enumerate(presentation.slides, start=1):
        if slide_index > MAX_SLIDES:
            truncated = True
            break
        parts = []
        title = clean_text(slide.shapes.title.text) if slide.shapes.title else ""
        if title:
            parts.append(f"Назва: {title}")
        read_shapes(slide.shapes, parts, title)
        try:
            notes_frame = slide.notes_slide.notes_text_frame
            notes = clean_text(notes_frame.text if notes_frame else "")
        except (AttributeError, ValueError):
            notes = ""
        if notes:
            parts.append("Нотатки доповідача:\n" + notes)
        if parts:
            slides.append(f"Слайд {slide_index}:\n" + "\n".join(parts))
        if sum(len(part) for part in slides) >= MAX_OUTPUT_CHARS:
            truncated = True
            break
    if len(presentation.slides) > MAX_SLIDES:
        truncated = True

    properties = presentation.core_properties
    content = clean_text("\n\n".join(slides))
    media_count = sum(1 for name in archive.namelist() if name.startswith("ppt/media/"))
    return {
        "content": content,
        "metadata": {
            "format": "pptx",
            "slides_total": len(presentation.slides),
            "slides_read": min(len(presentation.slides), MAX_SLIDES),
            "images": max(image_count, media_count),
            "tables": table_count,
            "charts": chart_count,
            "title": clean_text(properties.title),
            "author": clean_text(properties.author),
            "truncated": truncated,
        },
        "needs_ai_fallback": len(content) < 120 and max(image_count, media_count) > 0,
    }


EXTRACTORS = {
    ".docx": extract_docx,
    ".xlsx": extract_xlsx,
    ".pptx": extract_pptx,
}


def extract_document(payload: dict) -> dict:
    filename = Path(str(payload.get("filename") or "document")).name
    extension = Path(filename.lower()).suffix
    if extension not in SUPPORTED_EXTENSIONS:
        raise ValueError("unsupported_document_type")
    encoded = str(payload.get("data_base64") or "")
    try:
        data = base64.b64decode(encoded, validate=True)
    except (ValueError, TypeError) as exc:
        raise ValueError("invalid_base64") from exc
    if not data or len(data) > MAX_FILE_BYTES:
        raise ValueError("invalid_file_size")

    archive = validate_office_zip(data)
    try:
        result = EXTRACTORS[extension](data, archive)
    finally:
        archive.close()
    full_content = clean_text(result["content"])
    truncated = len(full_content) > MAX_OUTPUT_CHARS or bool(result["metadata"].get("truncated"))
    content = full_content[:MAX_OUTPUT_CHARS].strip()
    return {
        "ok": True,
        "filename": filename,
        "content": content,
        "characters": len(full_content),
        "stored_characters": len(content),
        "truncated": truncated,
        "needs_ai_fallback": bool(result["needs_ai_fallback"]),
        "metadata": result["metadata"],
    }


class Handler(BaseHTTPRequestHandler):
    server_version = "InboxDocumentReader/1.0"

    def log_message(self, fmt, *args):
        return

    def send_json(self, status: int, payload: dict):
        body = json.dumps(payload, ensure_ascii=False).encode("utf-8")
        self.send_response(status)
        self.send_header("Content-Type", "application/json; charset=utf-8")
        self.send_header("Content-Length", str(len(body)))
        self.end_headers()
        self.wfile.write(body)

    def do_GET(self):
        if self.path == "/healthz":
            self.send_json(200, {"status": "ok"})
        else:
            self.send_json(404, {"ok": False, "error": "not_found"})

    def do_POST(self):
        if self.path != "/extract":
            self.send_json(404, {"ok": False, "error": "not_found"})
            return
        try:
            length = int(self.headers.get("Content-Length", "0"))
            if length <= 0 or length > MAX_BODY_BYTES:
                raise ValueError("invalid_body_size")
            payload = json.loads(self.rfile.read(length))
            self.send_json(200, extract_document(payload))
        except Exception as exc:
            self.send_json(200, {"ok": False, "error": str(exc)[:1000]})


if __name__ == "__main__":
    ThreadingHTTPServer(("0.0.0.0", 8080), Handler).serve_forever()
